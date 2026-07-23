# Generates OpsGenie's investor pitch deck (pitch/OpsGenie_Pitch_Deck.pptx).
# Run with the standalone interpreter that has python-pptx installed:
#   C:\Python314\python.exe pitch\build_deck.py
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---- Brand palette (from app/static/css/dashboard.css) ----
DARK = RGBColor(0x06, 0x4e, 0x3b)      # --brand-dark
EMERALD = RGBColor(0x10, 0xb9, 0x81)   # --brand
MINT = RGBColor(0x34, 0xd3, 0x99)      # --brand-bright
BRIGHT_DARK = RGBColor(0x05, 0x96, 0x69)  # --brand-bright-dark
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x0f, 0x17, 0x2a)       # near-black body text
SLATE = RGBColor(0x47, 0x55, 0x69)     # secondary text
LIGHT_BG = RGBColor(0xf0, 0xfd, 0xf9)  # mint-tinted light panel
CARD_BORDER = RGBColor(0xd1, 0xf5, 0xe8)
AMBER = RGBColor(0xd9, 0x77, 0x06)     # for "ask" / caution placeholders

FONT = "Segoe UI"
FONT_SEMIBOLD = "Segoe UI Semibold"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def rect(slide, x, y, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if not line:
        no_line(shp)
    shp.shadow.inherit = False
    return shp


def rrect(slide, x, y, w, h, color, radius=0.06, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1)
    else:
        no_line(shp)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, valign=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def set_run(run, text, size, color, bold=False, italic=False, font=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def para(tf, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
          space_after=0, space_before=0, font=FONT, first=False, line_spacing=None):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    set_run(r, text, size, color, bold, italic, font)
    return p


def add_page_chrome(slide, page_no, section=""):
    # thin footer rule
    rect(slide, Inches(0.55), Inches(7.08), Inches(12.23), Pt(1.2), CARD_BORDER)
    tb, tf = textbox(slide, Inches(0.55), Inches(7.12), Inches(6), Inches(0.3))
    para(tf, "OpsGenie", 10, SLATE, bold=True, first=True)
    tb2, tf2 = textbox(slide, Inches(11.4), Inches(7.12), Inches(1.4), Inches(0.3))
    para(tf2, str(page_no), 10, SLATE, align=PP_ALIGN.RIGHT, first=True)


def add_header(slide, eyebrow, title, title_size=30, title_color=DARK):
    rect(slide, Inches(0.55), Inches(0.5), Inches(0.5), Inches(0.07), EMERALD)
    tb, tf = textbox(slide, Inches(0.55), Inches(0.62), Inches(11.5), Inches(0.35))
    para(tf, eyebrow.upper(), 12.5, EMERALD, bold=True, first=True)
    tb2, tf2 = textbox(slide, Inches(0.55), Inches(0.95), Inches(12.2), Inches(1.05))
    para(tf2, title, title_size, title_color, bold=True, first=True, line_spacing=1.0)


# ---------------------------------------------------------------------------
# 1. TITLE
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s, DARK)
rect(s, Inches(0), Inches(6.55), SLIDE_W, Inches(0.06), EMERALD)
# corner accent
c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(-1.6), Inches(4.2), Inches(4.2))
c.fill.solid(); c.fill.fore_color.rgb = BRIGHT_DARK; no_line(c); c.shadow.inherit = False
c.fill.fore_color.brightness = 0
sp = c.fill.fore_color._xFill
# transparency via alpha on solidFill
a = sp.find(qn('a:srgbClr'))
alpha = a.makeelement(qn('a:alpha'), {'val': '18000'})
a.append(alpha)

tb, tf = textbox(s, Inches(0.9), Inches(2.55), Inches(9), Inches(1.3))
para(tf, "OpsGenie", 60, WHITE, bold=True, first=True)
tb2, tf2 = textbox(s, Inches(0.95), Inches(3.65), Inches(10.8), Inches(0.8))
para(tf2, "The WhatsApp-first financial operating assistant for B2B distributors", 22, MINT, bold=True, first=True)
tb3, tf3 = textbox(s, Inches(0.95), Inches(4.35), Inches(10.5), Inches(0.7))
para(tf3, "Turns the records you already keep into a decision every morning — before the day begins.", 15, RGBColor(0xd1, 0xfa, 0xe5), first=True, italic=True)

tb4, tf4 = textbox(s, Inches(0.95), Inches(6.65), Inches(8), Inches(0.5))
para(tf4, "[Your Name]  ·  Founder  ·  tripathyspandan23@gmail.com", 12, RGBColor(0xa7, 0xf3, 0xd0), first=True)

# ---------------------------------------------------------------------------
# 2. THE PROBLEM
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s, WHITE)
add_header(s, "The Problem", "A distributor’s financial reality lives in Tally.\nTheir daily decisions happen on WhatsApp and phone calls.")
items = [
    ("Data ≠ decisions", "Books are accurate, but nobody turns them into “what do I do today” fast enough to matter."),
    ("Monthly cash surprises", "Distributors managing 20–50 dealer accounts hit cash-flow surprises monthly or more often — discovered too late to act on."),
    ("Zero appetite for new software", "They already run Tally/Vyapar and live on WhatsApp. A new app to learn is a new app that goes unused."),
]
x = Inches(0.55); w = Inches(3.95); gap = Inches(0.18)
for i, (h, d) in enumerate(items):
    cx = x + i * (w + gap)
    card = rrect(s, cx, Inches(2.35), w, Inches(3.6), LIGHT_BG, radius=0.06, line_color=CARD_BORDER)
    numtb, numtf = textbox(s, cx + Inches(0.3), Inches(2.6), Inches(1), Inches(0.6))
    para(numtf, f"0{i+1}", 26, MINT, bold=True, first=True)
    htb, htf = textbox(s, cx + Inches(0.3), Inches(3.25), w - Inches(0.6), Inches(0.8))
    para(htf, h, 17, DARK, bold=True, first=True, line_spacing=1.05)
    dtb, dtf = textbox(s, cx + Inches(0.3), Inches(3.95), w - Inches(0.6), Inches(1.8))
    para(dtf, d, 12.5, SLATE, first=True, line_spacing=1.25)
add_page_chrome(s, 2)

# ---------------------------------------------------------------------------
# 3. THE INSIGHT (statement slide)
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s, DARK)
rect(s, Inches(0.9), Inches(2.15), Inches(0.6), Inches(0.06), MINT)
tb, tf = textbox(s, Inches(0.9), Inches(2.4), Inches(11.2), Inches(2.6), valign=MSO_ANCHOR.TOP)
para(tf, "Distributors don’t get blindsided because they lack data.", 30, WHITE, bold=True, first=True, line_spacing=1.15)
para(tf, "They get blindsided because their data never becomes a decision fast enough.", 30, MINT, bold=True, line_spacing=1.15, space_before=6)
tb2, tf2 = textbox(s, Inches(0.9), Inches(5.15), Inches(10.5), Inches(0.9))
para(tf2, "OpsGenie closes that gap every single morning — automatically, on the channel they’re already in.", 16, RGBColor(0xd1, 0xfa, 0xe5), italic=True, first=True)
add_page_chrome(s, 3)

# ---------------------------------------------------------------------------
# 4. THE SOLUTION
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s, WHITE)
add_header(s, "The Solution", "OpsGenie — a deterministic financial engine, delivered entirely over WhatsApp")
tb, tf = textbox(s, Inches(0.55), Inches(2.15), Inches(6.0), Inches(4.4))
bullets = [
    "Ingests what a distributor already has — Tally, Vyapar, Excel exports, even photographed Tally voucher PDFs.",
    "Maintains one deterministic ledger of invoices, payments, dealers and suppliers — no re-entry, no parallel system of record.",
    "Delivers a daily cash-position briefing, on-demand reports, and guided write actions (record a payment, create an order, update stock) entirely on WhatsApp.",
    "An LLM narrates and answers free-form questions — but never owns business state and never does money math.",
]
for i, b in enumerate(bullets):
    p = para(tf, "", 13.5, INK, first=(i == 0), space_after=14, line_spacing=1.25)
    r = p.add_run(); set_run(r, "▸  ", 13.5, EMERALD, bold=True)
    r2 = p.add_run(); set_run(r2, b, 13.5, INK)

# right-side mini phone mockup illustrating the loop
px, py, pw, ph = Inches(7.05), Inches(2.15), Inches(5.7), Inches(4.5)
panel = rrect(s, px, py, pw, ph, DARK, radius=0.05)
tb2, tf2 = textbox(s, px + Inches(0.35), py + Inches(0.28), pw - Inches(0.7), Inches(0.4))
para(tf2, "8:00 AM — Morning Briefing", 13, MINT, bold=True, first=True)
lines = [
    "Good morning! Cash position: ₹4.2L in, ₹1.8L out (7d)",
    "⚠️ Watch this week: Sharma Traders — ₹78,000 due in 2 days",
    "📦 3 products may run out within 10 days",
    "Reply with a number, or type your question.",
]
yy = py + Inches(0.85)
for ln in lines:
    bubble = rrect(s, px + Inches(0.35), yy, pw - Inches(0.9), Inches(0.62), RGBColor(0x0b, 0x3a, 0x2c), radius=0.25)
    btb, btf = textbox(s, px + Inches(0.55), yy + Inches(0.08), pw - Inches(1.3), Inches(0.5), valign=MSO_ANCHOR.MIDDLE)
    para(btf, ln, 10.8, RGBColor(0xea, 0xfb, 0xf3), first=True, line_spacing=1.1)
    yy += Inches(0.78)
capt, captf = textbox(s, px + Inches(0.35), py + ph - Inches(0.55), pw - Inches(0.7), Inches(0.4))
para(captf, "See the full working demo → pitch/demo/opsgenie_demo.html", 10, MINT, italic=True, first=True)
add_page_chrome(s, 4)

# ---------------------------------------------------------------------------
# 5. THE PRODUCT LOOP
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s, WHITE)
add_header(s, "How It Works", "The product loop — this is the moat, not any single feature")
steps = [
    "Business event occurs", "Deterministic engine processes it", "Snapshot reflects new reality",
    "Recommendation engine runs", "Morning briefing delivered", "Distributor takes action",
]
n = len(steps)
total_w = Inches(12.2)
box_w = Inches(1.72)
arrow_w = Emu(int((total_w - box_w * n) / (n - 1)))
xcur = Inches(0.55)
ytop = Inches(2.5)
for i, step in enumerate(steps):
    col = DARK if i % 2 == 0 else BRIGHT_DARK
    b = rrect(s, xcur, ytop, box_w, Inches(1.5), col, radius=0.1)
    btb, btf = textbox(s, xcur + Inches(0.12), ytop + Inches(0.15), box_w - Inches(0.24), Inches(1.2), valign=MSO_ANCHOR.MIDDLE)
    para(btf, step, 11.5, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, line_spacing=1.05)
    xcur += box_w
    if i < n - 1:
        atb, atf = textbox(s, xcur, ytop + Inches(0.45), arrow_w, Inches(0.6), valign=MSO_ANCHOR.MIDDLE)
        para(atf, "→", 22, MINT, bold=True, align=PP_ALIGN.CENTER, first=True)
        xcur += arrow_w
loop_tb, loop_tf = textbox(s, Inches(0.55), Inches(4.5), Inches(12.2), Inches(0.5))
para(loop_tf, "… which produces a new business event, and the loop repeats — daily, compounding.", 13, SLATE, italic=True, align=PP_ALIGN.CENTER, first=True)

quote_panel = rrect(s, Inches(1.4), Inches(5.35), Inches(10.5), Inches(1.35), LIGHT_BG, radius=0.08, line_color=CARD_BORDER)
qtb, qtf = textbox(s, Inches(1.75), Inches(5.55), Inches(9.8), Inches(1.0), valign=MSO_ANCHOR.MIDDLE)
para(qtf, "“The morning briefing is the output of this loop, not the product itself. "
          "The moat is the loop running daily, with real data accumulating over months.”", 14.5, DARK, italic=True, bold=True, align=PP_ALIGN.CENTER, first=True, line_spacing=1.2)
add_page_chrome(s, 5)

# ---------------------------------------------------------------------------
# 6. TRUST / ARCHITECTURE PRINCIPLE
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s, WHITE)
add_header(s, "Why Distributors Can Trust It", "Money math never touches the LLM — by construction, not by policy")
principles = [
    ("Deterministic ledger", "Every number in every briefing traces back to a real invoice or payment. The LLM narrates; it never calculates, decides, or writes."),
    ("Writes are guided workflows", "Creating an order or recording a payment is a confirm-gated flow (“Reply YES to confirm”), re-derived fresh at the moment of confirmation — never a freeform AI action."),
    ("Money-safety guard on free-form Q&A", "The AI Q&A agent discards any reply containing a figure its own read-only tools didn’t actually return."),
    ("Provider-agnostic & resilient", "A pluggable failover chain across 6 LLM providers (Claude, Gemini, Groq, OpenRouter, GitHub Models, Cohere) — verified in production failing over a real 429 with identical figures."),
]
gy = Inches(2.15)
for i, (h, d) in enumerate(principles):
    row = i // 2
    col = i % 2
    cx = Inches(0.55) + col * Inches(6.15)
    cy = gy + row * Inches(2.4)
    rrect(s, cx, cy, Inches(5.95), Inches(2.2), LIGHT_BG, radius=0.06, line_color=CARD_BORDER)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.3), cy + Inches(0.3), Inches(0.16), Inches(0.16))
    dot.fill.solid(); dot.fill.fore_color.rgb = EMERALD; no_line(dot); dot.shadow.inherit = False
    htb, htf = textbox(s, cx + Inches(0.3), cy + Inches(0.58), Inches(5.35), Inches(0.5))
    para(htf, h, 15.5, DARK, bold=True, first=True)
    dtb, dtf = textbox(s, cx + Inches(0.3), cy + Inches(1.05), Inches(5.35), Inches(1.05))
    para(dtf, d, 11.8, SLATE, first=True, line_spacing=1.22)
add_page_chrome(s, 6)


def feature_grid_slide(page_no, eyebrow, title, cards):
    s = add_slide()
    set_bg(s, WHITE)
    add_header(s, eyebrow, title, title_size=26)
    cols = 3
    cw = Inches(3.95); ch = Inches(2.15); gapx = Inches(0.18); gapy = Inches(0.2)
    x0 = Inches(0.55); y0 = Inches(2.05)
    for i, (h, d) in enumerate(cards):
        r_, c_ = divmod(i, cols)
        cx = x0 + c_ * (cw + gapx)
        cy = y0 + r_ * (ch + gapy)
        rrect(s, cx, cy, cw, ch, LIGHT_BG, radius=0.07, line_color=CARD_BORDER)
        rect(s, cx, cy, Inches(0.07), ch, EMERALD)
        htb, htf = textbox(s, cx + Inches(0.28), cy + Inches(0.2), cw - Inches(0.5), Inches(0.55))
        para(htf, h, 13.5, DARK, bold=True, first=True, line_spacing=1.05)
        dtb, dtf = textbox(s, cx + Inches(0.28), cy + Inches(0.75), cw - Inches(0.5), ch - Inches(0.95))
        para(dtf, d, 10.8, SLATE, first=True, line_spacing=1.2)
    add_page_chrome(s, page_no)
    return s


# ---------------------------------------------------------------------------
# 7. FEATURES — Daily Intelligence
# ---------------------------------------------------------------------------
feature_grid_slide(7, "Everything We’ve Built · 1 of 4", "Daily intelligence, delivered before the day starts", [
    ("Morning & evening briefings", "LLM-narrated cash position, honest accrual vs. cash-basis split — never blended into one number."),
    ("7-day cash-shortage forecast", "Walks expected collections/payments day-by-day, names the single biggest payment driving the shortfall."),
    ("Stock-out forecast", "Sales-velocity model per product, guarded against thin data (≥3 sale-days, ≥10 units) before it ever speaks up."),
    ("Pre-due invoice nudges", "Fires 1–2 days before a receivable is due — with a suggested action, not just a number."),
    ("Dealer-risk & overdue alerts", "Automatic flags for high-risk dealers; optional direct-to-dealer reminders with the distributor’s consent."),
    ("Trend analytics", "Week-over-week cash trends and 30-day dealer/product movement — including credit-risk dealers whose balance is climbing."),
])

# ---------------------------------------------------------------------------
# 8. FEATURES — WhatsApp-native writes
# ---------------------------------------------------------------------------
feature_grid_slide(8, "Everything We’ve Built · 2 of 4", "A real system of record — written entirely from WhatsApp", [
    ("Create orders & invoices", "Guided flow with live GST math, credit-limit and duplicate-order warnings, and an optional advance-payment step."),
    ("Record payments (FIFO)", "Same allocator used by CSV import — confirm-gated, re-derived fresh at the moment of commit."),
    ("Invoice photo → OCR pre-fill", "Photograph a paper invoice; vision-LLM pre-fills the same guided flow field-by-field — every field still needs a human YES."),
    ("Branded PDF invoices", "Auto-generated, Unicode-script aware (Hindi/Odia/English), delivered straight to the dealer on WhatsApp."),
    ("Edit, void & stock-take", "Full audit trail (old → new, optional reason) on every correction; risky auto-reallocation deliberately left out for safety."),
    ("Marketing broadcast", "Opt-in only, segment by all/overdue/named dealers, re-derives the recipient list fresh at send time."),
])

# ---------------------------------------------------------------------------
# 9. FEATURES — Reports & data
# ---------------------------------------------------------------------------
feature_grid_slide(9, "Everything We’ve Built · 3 of 4", "Vyapar/Tally-grade reports, without leaving WhatsApp", [
    ("GST sales & purchase registers", "Rate-wise summary, per-invoice-line detail — requested straight from a WhatsApp chat."),
    ("Ledger, day book & aging report", "Opening → running → closing balance; aging bucketed Not Due / 0-30 / 31-60 / 61-90 / 90+."),
    ("Full Excel export", "An 11-sheet workbook per company, via a short-lived signed link — no login required."),
    ("Flexible import pipeline", "Tally, Vyapar, canonical CSV/Excel — and now Tally “Multi Voucher Print” PDFs, parsed directly."),
    ("Self-serve reconciliation", "Import totals shown and confirmed before WhatsApp ever starts — catch mistakes on a form, not mid-conversation."),
    ("Delivery status visibility", "“Was that invoice PDF actually delivered?” — answerable on demand from real Meta delivery receipts."),
])

# ---------------------------------------------------------------------------
# 10. FEATURES — Reach & trust
# ---------------------------------------------------------------------------
feature_grid_slide(10, "Everything We’ve Built · 4 of 4", "Built for how India’s distributors actually operate", [
    ("5-locale multilingual", "English, Hindi & Odia — Devanagari/Odia script and Romanized-first variants, tuned for how people actually type on WhatsApp."),
    ("Unicode invoice PDFs", "Correctly renders regional-script dealer/product names — not just Latin text."),
    ("Self-serve onboarding", "Public web wizard — import your data or start fresh, then WhatsApp only asks for what wasn’t already provided."),
    ("Admin dashboard", "Password-gated, founder-facing view across every company, invoice, and payment."),
    ("Security by design", "HMAC-verified WhatsApp webhook, signed short-lived export links, capability-token-gated onboarding."),
    ("Resumable, correctable setup", "Progress checklist, restart, and 30+ guided correction commands — nothing is a one-shot, unrecoverable step."),
])

# ---------------------------------------------------------------------------
# 11. UNDER THE HOOD
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s, DARK)
add_header(s, "Under the Hood", "Engineered like production software, not a demo", title_color=WHITE)
stats = [
    ("17", "ORM tables"), ("44", "Alembic migrations"), ("69", "test files"),
    ("1,002", "automated tests"), ("6", "LLM providers, auto-failover"), ("5", "languages/scripts supported"),
]
cw = Inches(3.9); ch = Inches(1.55); gapx = Inches(0.2); gapy = Inches(0.22)
x0 = Inches(0.6); y0 = Inches(2.15)
for i, (num, label) in enumerate(stats):
    r_, c_ = divmod(i, 3)
    cx = x0 + c_ * (cw + gapx)
    cy = y0 + r_ * (ch + gapy)
    rrect(s, cx, cy, cw, ch, RGBColor(0x0b, 0x3a, 0x2c), radius=0.1)
    ntb, ntf = textbox(s, cx + Inches(0.25), cy + Inches(0.18), cw - Inches(0.5), Inches(0.7))
    para(ntf, num, 30, MINT, bold=True, first=True)
    ltb, ltf = textbox(s, cx + Inches(0.27), cy + Inches(0.98), cw - Inches(0.5), Inches(0.5))
    para(ltf, label, 12, RGBColor(0xd1, 0xfa, 0xe5), first=True)
stack_tb, stack_tf = textbox(s, Inches(0.6), Inches(5.85), Inches(11.9), Inches(1.0))
para(stack_tf, "FastAPI · PostgreSQL (Neon) · Meta WhatsApp Business API · Render (API + scheduler) · Vercel (marketing site) · fpdf2 · APScheduler", 12.5, MINT, align=PP_ALIGN.CENTER, first=True)
add_page_chrome(s, 11)

# ---------------------------------------------------------------------------
# 12. TRACTION
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s, WHITE)
add_header(s, "Traction", "Built solo. Shipped in production. Running on a real pilot.")
left_items = [
    "Live in production for a real distributor pilot on WhatsApp — not a sandboxed demo.",
    "Every roadmap version shipped: V0.0 Proof of Value → V0.1 Operational Product → V0.2 Source of Truth → V0.3 Intelligence & Expansion — all complete.",
    "A full whole-codebase audit (≈32k LOC, 4 parallel domain reviews) found and fixed 8 real production-grade bugs — each shipped with a regression test proven to fail without the fix.",
    "Live incidents diagnosed and resolved directly from the founder’s own WhatsApp screenshots while distributors were actively using it.",
]
tb, tf = textbox(s, Inches(0.55), Inches(2.0), Inches(7.3), Inches(4.6))
for i, b in enumerate(left_items):
    p = para(tf, "", 13, INK, first=(i == 0), space_after=16, line_spacing=1.28)
    r = p.add_run(); set_run(r, "✓  ", 13, EMERALD, bold=True)
    r2 = p.add_run(); set_run(r2, b, 13, INK)

panel = rrect(s, Inches(8.15), Inches(2.0), Inches(4.6), Inches(4.55), LIGHT_BG, radius=0.06, line_color=CARD_BORDER)
ptb, ptf = textbox(s, Inches(8.45), Inches(2.25), Inches(4.0), Inches(0.4))
para(ptf, "Version roadmap", 13, DARK, bold=True, first=True)
versions = [
    ("V0.0", "Proof of Value", True), ("V0.1", "Operational Product", True),
    ("V0.2", "Source of Truth", True), ("V0.3", "Intelligence & Expansion", True),
    ("Post-SPEC", "Onboarding, dashboard, multilingual, OCR, forecasts", True),
]
yy = Inches(2.75)
for tag, label, done in versions:
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.45), yy + Inches(0.04), Inches(0.16), Inches(0.16))
    dot.fill.solid(); dot.fill.fore_color.rgb = EMERALD if done else SLATE; no_line(dot); dot.shadow.inherit = False
    vtb, vtf = textbox(s, Inches(8.75), yy - Inches(0.03), Inches(3.9), Inches(0.55))
    p = para(vtf, "", 11.5, DARK, first=True, line_spacing=1.1)
    r = p.add_run(); set_run(r, f"{tag}  ", 11.5, EMERALD, bold=True)
    r2 = p.add_run(); set_run(r2, label, 11.5, INK)
    yy += Inches(0.72)
add_page_chrome(s, 12)

# ---------------------------------------------------------------------------
# 13. MARKET
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s, WHITE)
add_header(s, "Market", "A large, underserved, WhatsApp-native customer base")
tb, tf = textbox(s, Inches(0.55), Inches(2.0), Inches(6.6), Inches(4.6))
paras = [
    ("Who we serve", "FMCG, pharma, poultry/agri and wholesale B2B distributors already running Tally or Vyapar, managing 20–50 dealer accounts, facing cash-flow surprises monthly or more often."),
    ("Why WhatsApp is the wedge", "It's already their primary business channel — no app to install, no login to remember, no training session required."),
    ("Why now", "WhatsApp Business API access, capable low-cost LLMs, and India's UPI/GST-driven digitization of small-business records have only recently made this combination viable."),
]
for i, (h, d) in enumerate(paras):
    para(tf, h, 15, DARK, bold=True, first=(i == 0), space_before=(0 if i == 0 else 16), space_after=4)
    para(tf, d, 12.5, SLATE, space_after=0, line_spacing=1.25)

panel = rrect(s, Inches(7.55), Inches(2.0), Inches(5.2), Inches(4.55), DARK, radius=0.06)
ptb, ptf = textbox(s, Inches(7.85), Inches(2.25), Inches(4.6), Inches(0.4))
para(ptf, "Sizing (illustrative — refine before pitching)", 12, MINT, bold=True, first=True)
rows = [
    ("Tally-using businesses in India", "2M+ (public Tally figure)"),
    ("Est. B2B distributors, target profile", "[X hundred thousand — size this]"),
    ("SAM (reachable via WhatsApp, GST-registered)", "[refine]"),
    ("SOM (Yr 1 pilot geography)", "[refine]"),
]
yy = Inches(2.85)
for label, val in rows:
    ltb, ltf = textbox(s, Inches(7.85), yy, Inches(4.6), Inches(0.4))
    para(ltf, label, 11, RGBColor(0xd1, 0xfa, 0xe5), first=True)
    vtb, vtf = textbox(s, Inches(7.85), yy + Inches(0.32), Inches(4.6), Inches(0.4))
    col = AMBER if "[" in val else MINT
    para(vtf, val, 13.5, col, bold=True, first=True)
    yy += Inches(0.9)
add_page_chrome(s, 13)

# ---------------------------------------------------------------------------
# 14. BUSINESS MODEL
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s, WHITE)
add_header(s, "Business Model", "Simple SaaS, distribution through the network already in the chat")
cards = [
    ("Subscription", "₹999/month per distributor (pilot pricing). Founder-activated today; self-serve payment is the next infrastructure step."),
    ("Low CAC channel", "WhatsApp is both the product and the distribution surface — a distributor’s own dealer/supplier network sees the product working."),
    ("Expansion path", "Land with one distributor → their accountant/operator becomes the champion → word-of-mouth within tight-knit distributor trade circles."),
]
x = Inches(0.55); w = Inches(3.95); gap = Inches(0.18)
for i, (h, d) in enumerate(cards):
    cx = x + i * (w + gap)
    rrect(s, cx, Inches(2.15), w, Inches(2.9), LIGHT_BG, radius=0.07, line_color=CARD_BORDER)
    htb, htf = textbox(s, cx + Inches(0.3), Inches(2.4), w - Inches(0.6), Inches(0.6))
    para(htf, h, 15.5, DARK, bold=True, first=True)
    dtb, dtf = textbox(s, cx + Inches(0.3), Inches(3.0), w - Inches(0.6), Inches(1.9))
    para(dtf, d, 12, SLATE, first=True, line_spacing=1.25)
note_tb, note_tf = textbox(s, Inches(0.55), Inches(5.4), Inches(12.2), Inches(1.2))
para(note_tf, "Known, deliberately-deferred gap: no payment gateway yet — subscription activation is manual today (correct sequencing for pilot validation, per this project’s phase-discipline principle: prove the product before automating billing for it).", 11.5, SLATE, italic=True, first=True, line_spacing=1.25)
add_page_chrome(s, 14)

# ---------------------------------------------------------------------------
# 15. ROADMAP
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s, WHITE)
add_header(s, "What’s Next", "Harden for scale, then widen the funnel")
cols = [
    ("Harden for scale", ["CI pipeline running the full suite on every change", "LLM provider timeouts + scheduler catch-up", "Security hardening (rate limiting, session expiry)"]),
    ("Widen the product", ["Voice notes (low-literacy operators)", "Self-serve payment gateway", "GSTR-filing-ready export tier"]),
    ("Grow the pilot base", ["Formalize the current pilot into a case study", "Onboard the next cohort of distributors", "Meta-approved marketing/reminder templates"]),
]
x = Inches(0.55); w = Inches(3.95); gap = Inches(0.18)
for i, (h, items) in enumerate(cols):
    cx = x + i * (w + gap)
    rrect(s, cx, Inches(2.05), w, Inches(4.3), LIGHT_BG, radius=0.07, line_color=CARD_BORDER)
    rect(s, cx, Inches(2.05), w, Inches(0.62), DARK)
    htb, htf = textbox(s, cx + Inches(0.25), Inches(2.18), w - Inches(0.5), Inches(0.4), valign=MSO_ANCHOR.MIDDLE)
    para(htf, h, 13.5, WHITE, bold=True, first=True)
    ytb, ytf = textbox(s, cx + Inches(0.3), Inches(2.9), w - Inches(0.55), Inches(3.3))
    for j, it in enumerate(items):
        p = para(ytf, "", 11.8, INK, first=(j == 0), space_after=12, line_spacing=1.25)
        r = p.add_run(); set_run(r, "•  ", 11.8, EMERALD, bold=True)
        r2 = p.add_run(); set_run(r2, it, 11.8, INK)
add_page_chrome(s, 15)

# ---------------------------------------------------------------------------
# 16. THE ASK
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s, DARK)
add_header(s, "The Ask", "[Placeholder — fill in before presenting]", title_color=WHITE)
tb, tf = textbox(s, Inches(0.55), Inches(2.1), Inches(11.9), Inches(1.0))
para(tf, "Raising ₹[X] to fund the next [12/18] months", 24, MINT, bold=True, first=True)
uses = [
    "[N] engineering hires to build CI, harden scheduling/security, and ship the next feature tier",
    "Go-to-market into [N] additional pilot distributors across [geography/vertical]",
    "Meta template approvals + WhatsApp Business API scale-up costs",
    "Runway to a self-serve payment gateway and repeatable onboarding motion",
]
tb2, tf2 = textbox(s, Inches(0.55), Inches(3.15), Inches(11.5), Inches(3.0))
for i, u in enumerate(uses):
    p = para(tf2, "", 14, WHITE, first=(i == 0), space_after=14, line_spacing=1.3)
    r = p.add_run(); set_run(r, "▸  ", 14, MINT, bold=True)
    r2 = p.add_run(); set_run(r2, u, 14, RGBColor(0xea, 0xfb, 0xf3))
note_tb, note_tf = textbox(s, Inches(0.55), Inches(6.35), Inches(11.5), Inches(0.6))
para(note_tf, "Note: figures on this slide are placeholders — fill in your actual raise amount, hiring plan, and geography before sending this deck out.", 10.5, AMBER, italic=True, first=True)
add_page_chrome(s, 16)

# ---------------------------------------------------------------------------
# 17. CLOSING
# ---------------------------------------------------------------------------
s = add_slide()
set_bg(s, DARK)
c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(4.2), Inches(4.2), Inches(4.2))
c.fill.solid(); c.fill.fore_color.rgb = BRIGHT_DARK; no_line(c); c.shadow.inherit = False
a2 = c.fill.fore_color._xFill.find(qn('a:srgbClr'))
alpha2 = a2.makeelement(qn('a:alpha'), {'val': '18000'}); a2.append(alpha2)
tb, tf = textbox(s, Inches(0.9), Inches(2.9), Inches(10), Inches(1.0))
para(tf, "Every morning, before the day begins.", 32, WHITE, bold=True, first=True)
tb2, tf2 = textbox(s, Inches(0.95), Inches(3.75), Inches(10), Inches(0.6))
para(tf2, "OpsGenie · tripathyspandan23@gmail.com", 15, MINT, first=True)

out_path = "OpsGenie_Pitch_Deck.pptx"
prs.save(out_path)
print(f"Saved {out_path} with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
