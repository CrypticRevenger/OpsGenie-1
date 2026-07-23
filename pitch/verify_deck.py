# Static overflow checker for OpsGenie_Pitch_Deck.pptx.
# No LibreOffice/PowerPoint available to render locally, so this re-measures
# each textbox's wrapped text with the real Segoe UI font metrics and flags:
#   - shapes that fall outside the slide canvas
#   - text blocks whose wrapped height exceeds their box's allocated height
from pptx import Presentation
from pptx.util import Emu
from PIL import ImageFont

EMU_PER_IN = 914400
FONT_DIR = "C:/Windows/Fonts/"
FONT_FILES = {
    (False, False): FONT_DIR + "segoeui.ttf",
    (True, False): FONT_DIR + "segoeuib.ttf",
    (False, True): FONT_DIR + "segoeuii.ttf",
    (True, True): FONT_DIR + "segoeuiz.ttf",
}
_font_cache = {}


def get_font(size_pt, bold, italic):
    key = (round(size_pt), bold, italic)
    if key not in _font_cache:
        path = FONT_FILES.get((bold, italic), FONT_FILES[(False, False)])
        # size in px at 72dpi == size in pt
        _font_cache[key] = ImageFont.truetype(path, max(1, round(size_pt)))
    return _font_cache[key]


def wrap_and_measure(text, size_pt, bold, italic, max_width_in):
    font = get_font(size_pt, bold, italic)
    max_width_px = max_width_in * 72  # 72dpi == pt-equals-px assumption
    words = text.split(" ")
    lines = 1
    cur = ""
    for w in words:
        trial = (cur + " " + w).strip()
        width = font.getlength(trial)
        if width > max_width_px and cur:
            lines += 1
            cur = w
        else:
            cur = trial
    return lines


def emu_in(v):
    return v / EMU_PER_IN


prs = Presentation("OpsGenie_Pitch_Deck.pptx")
slide_w_in = emu_in(prs.slide_width)
slide_h_in = emu_in(prs.slide_height)

issues = []

for si, slide in enumerate(slide.slides if False else prs.slides, start=1):
    for shape in slide.shapes:
        if shape.left is None:
            continue
        l, t, w, h = emu_in(shape.left), emu_in(shape.top), emu_in(shape.width), emu_in(shape.height)
        # bounds check (allow tiny negative for decorative off-canvas circles)
        if l < -2 or t < -2 or (l + w) > slide_w_in + 0.05 or (t + h) > slide_h_in + 0.05:
            # decorative corner ovals are intentionally off-canvas; skip those (no text frame or huge ovals)
            if not shape.has_text_frame:
                continue
            issues.append(f"[Slide {si}] shape '{shape.shape_id}' bounds ({l:.2f},{t:.2f},{w:.2f},{h:.2f}) exceed canvas {slide_w_in:.2f}x{slide_h_in:.2f}")

        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        total_height_pt = 0.0
        any_text = False
        for p in tf.paragraphs:
            text = "".join(r.text for r in p.runs)
            if not text.strip():
                continue
            any_text = True
            sizes = [r.font.size.pt if r.font.size else 12 for r in p.runs]
            bolds = [bool(r.font.bold) for r in p.runs]
            italics = [bool(r.font.italic) for r in p.runs]
            size_pt = max(sizes)
            bold = bolds[0]
            italic = italics[0]
            lines = wrap_and_measure(text, size_pt, bold, italic, w)
            line_spacing = p.line_spacing if p.line_spacing else 1.0
            line_h_pt = size_pt * 1.22 * (line_spacing if isinstance(line_spacing, float) else 1.0)
            block_h_pt = lines * line_h_pt
            block_h_pt += (p.space_before.pt if p.space_before else 0)
            block_h_pt += (p.space_after.pt if p.space_after else 0)
            total_height_pt += block_h_pt
        if not any_text:
            continue
        total_height_in = total_height_pt / 72
        budget_in = h
        if total_height_in > budget_in + 0.05:
            preview = tf.paragraphs[0].runs[0].text[:40] if tf.paragraphs[0].runs else ""
            issues.append(
                f"[Slide {si}] textbox @({l:.2f},{t:.2f}) sized {w:.2f}x{h:.2f}in "
                f"needs ~{total_height_in:.2f}in for text ('{preview}...') — OVERFLOW by {total_height_in - budget_in:.2f}in"
            )

if issues:
    print(f"{len(issues)} potential issue(s):\n")
    for i in issues:
        print(" -", i)
else:
    print("No overflow issues detected.")
